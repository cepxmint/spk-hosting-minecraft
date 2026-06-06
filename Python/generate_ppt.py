import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 1. BACA DATASET & HITUNG TOPSIS
# ============================================================
df = pd.read_csv("Dokumen/Dataset_Hosting_Minecraft.csv")

kriteria = ["Harga (USD/mo)", "RAM (GB)", "Storage (GB)", "Slot Pemain", "Uptime (%)", "vCPU"]
jenis = ["cost", "benefit", "benefit", "benefit", "benefit", "benefit"]
bobot = [0.20, 0.20, 0.10, 0.10, 0.25, 0.15]

X = df[kriteria].values.astype(float)
R = X / np.sqrt((X ** 2).sum(axis=0))
Y = R * np.array(bobot)
A_pos = np.where(np.array(jenis) == "benefit", Y.max(axis=0), Y.min(axis=0))
A_neg = np.where(np.array(jenis) == "benefit", Y.min(axis=0), Y.max(axis=0))
D_pos = np.sqrt(((Y - A_pos) ** 2).sum(axis=1))
D_neg = np.sqrt(((Y - A_neg) ** 2).sum(axis=1))
V = D_neg / (D_pos + D_neg)
df["V"] = V
df["Rank"] = df["V"].rank(ascending=False).astype(int)
df_sorted = df.sort_values("Rank").reset_index(drop=True)

# ============================================================
# 2. SETUP PRESENTASI
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# Color palette
BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
GOLD = RGBColor(0xFF, 0xA5, 0x00)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_para(tf, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, space_after=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(space_after)
    return p


def make_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table


def style_header(table, texts, font_size=12):
    for j, text in enumerate(texts):
        cell = table.cell(0, j)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def style_cell(table, row, col, text, font_size=11, bold=False, fill=None, align=PP_ALIGN.LEFT):
    cell = table.cell(row, col)
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


# ============================================================
# SLIDE 1 — COVER
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, BLUE)

add_textbox(slide1, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
            "SISTEM PENDUKUNG KEPUTUSAN", 16, RGBColor(0xAA, 0xCC, 0xFF),
            bold=False, alignment=PP_ALIGN.CENTER)

add_textbox(slide1, Inches(1), Inches(2.3), Inches(11), Inches(1.5),
            "REKOMENDASI PLATFORM\nHOSTING SERVER MINECRAFT", 44, WHITE,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide1, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
            "Metode TOPSIS — 72 Alternatif | 6 Kriteria", 20, RGBColor(0xDD, 0xDD, 0xFF),
            alignment=PP_ALIGN.CENTER)

# Decorative line
add_shape(slide1, Inches(4.5), Inches(4), Inches(4.3), Inches(0.04), WHITE)

add_textbox(slide1, Inches(1), Inches(5.2), Inches(11), Inches(1),
            "Disusun oleh: Affandi Rahman Hakim — 23080960114", 16, RGBColor(0xCC, 0xDD, 0xFF),
            alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — DAFTAR ISI
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide2, Inches(0), Inches(0), Inches(0.15), H, BLUE)

add_textbox(slide2, Inches(1), Inches(0.5), Inches(5), Inches(0.8),
            "DAFTAR ISI", 36, BLUE, bold=True)

items = [
    ("01", "Latar Belakang & Studi Kasus"),
    ("02", "Dataset & Statistik"),
    ("03", "Kriteria & Pembobotan"),
    ("04", "Sekilas Metode TOPSIS"),
    ("05", "Langkah Perhitungan"),
    ("06", "Hasil & Analisis Ranking"),
    ("07", "Kesimpulan"),
]

for i, (num, title) in enumerate(items):
    y = 1.6 + i * 0.75
    add_textbox(slide2, Inches(1.2), Inches(y), Inches(0.6), Inches(0.5),
                num, 22, BLUE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_shape(slide2, Inches(1.9), Inches(y + 0.15), Inches(0.03), Inches(0.25), BLUE)
    add_textbox(slide2, Inches(2.1), Inches(y), Inches(9), Inches(0.5),
                title, 22, DARK)

# ============================================================
# SLIDE 3 — LATAR BELAKANG
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide3, Inches(0), Inches(0), Inches(0.15), H, BLUE)

add_textbox(slide3, Inches(1), Inches(0.5), Inches(8), Inches(0.8),
            "LATAR BELAKANG & STUDI KASUS", 36, BLUE, bold=True)

add_textbox(slide3, Inches(1), Inches(1.5), Inches(11), Inches(2.5),
            "Memilih platform hosting server Minecraft itu ribet! "
            "Ada banyak penyedia dengan harga, spesifikasi, dan kualitas "
            "yang berbeda-beda. Susah mana yang paling worth it.\n\n"
            "Anggap aja kayak milih kosan:\n"
            "  Harga = biaya sewa bulanan\n"
            "  RAM = luas kamar (makin gede makin lega)\n"
            "  Storage = ukuran lemari (nyimpen barang)\n"
            "  Slot Pemain = jumlah tamu yang bisa nginep\n"
            "  Uptime = keamanan & kenyamanan (jarang trouble)\n"
            "  vCPU = kecepatan WiFi\n\n"
            "Mana kosan terbaik? Biar ga pusing, kita pake TOPSIS!",
            22, DARK)

add_textbox(slide3, Inches(1), Inches(5.5), Inches(11), Inches(1.5),
            "Data diambil dari 22 platform hosting (72 tier plan) yang "
            "diverifikasi langsung dari website resmi dan artikel "
            "perbandingan pihak ketiga.",
            18, GRAY)

# ============================================================
# SLIDE 4 — DATASET & STATISTIK
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide4, Inches(0), Inches(0), Inches(0.15), H, BLUE)

add_textbox(slide4, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
            "DATASET & STATISTIK", 36, BLUE, bold=True)

# Summary boxes
summary = [
    ("72", "Alternatif"),
    ("22", "Platform"),
    ("6", "Kriteria"),
    ("$0 - $80", "Rentang Harga"),
    ("1 - 32 GB", "Rentang RAM"),
    ("99 - 100%", "Rentang Uptime"),
]

for i, (val, label) in enumerate(summary):
    col = i % 3
    row = i // 3
    x = 1.0 + col * 4.0
    y = 1.3 + row * 1.7
    box = add_shape(slide4, Inches(x), Inches(y), Inches(3.5), Inches(1.3), LIGHT_BLUE)
    add_textbox(slide4, Inches(x + 0.2), Inches(y + 0.1), Inches(3.1), Inches(0.6),
                val, 32, BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide4, Inches(x + 0.2), Inches(y + 0.6), Inches(3.1), Inches(0.5),
                label, 16, DARK, alignment=PP_ALIGN.CENTER)

# Sample data table
sample = df.head(5)
tbl = make_table(slide4, 6, 5, Inches(1), Inches(5), Inches(11), Inches(2.2))
style_header(tbl, ["No", "Platform", "Tier", "Harga", "RAM"])
for i in range(5):
    style_cell(tbl, i + 1, 0, str(i + 1), align=PP_ALIGN.CENTER)
    style_cell(tbl, i + 1, 1, sample.iloc[i]["Platform"])
    style_cell(tbl, i + 1, 2, sample.iloc[i]["Tier"])
    style_cell(tbl, i + 1, 3, f"${sample.iloc[i]['Harga (USD/mo)']:.2f}", align=PP_ALIGN.CENTER)
    style_cell(tbl, i + 1, 4, f"{sample.iloc[i]['RAM (GB)']:.0f} GB", align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5 — KRITERIA & BOBOT
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide5, Inches(0), Inches(0), Inches(0.15), H, BLUE)

add_textbox(slide5, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
            "KRITERIA & PEMBOBOTAN", 36, BLUE, bold=True)

criteria_data = [
    ("Harga (USD/mo)", "Cost", "0.20", "Makin murah makin baik"),
    ("RAM (GB)", "Benefit", "0.20", "Makin besar makin baik"),
    ("Storage (GB)", "Benefit", "0.10", "Makin besar makin baik"),
    ("Slot Pemain", "Benefit", "0.10", "Makin banyak makin baik"),
    ("Uptime (%)", "Benefit", "0.25", "Makin tinggi makin baik"),
    ("vCPU", "Benefit", "0.15", "Makin banyak makin baik"),
]

tbl = make_table(slide5, len(criteria_data) + 1, 4, Inches(0.8), Inches(1.4), Inches(11.5), Inches(3.5))
style_header(tbl, ["Kriteria", "Jenis", "Bobot", "Keterangan"], font_size=14)

for i, (krit, jns, bbt, ket) in enumerate(criteria_data):
    style_cell(tbl, i + 1, 0, krit, font_size=13, bold=True)
    style_cell(tbl, i + 1, 1, jns, font_size=13, align=PP_ALIGN.CENTER)
    style_cell(tbl, i + 1, 2, bbt, font_size=13, align=PP_ALIGN.CENTER)
    style_cell(tbl, i + 1, 3, ket, font_size=12, fill=LIGHT_BLUE)

add_textbox(slide5, Inches(1), Inches(5.5), Inches(11), Inches(1.5),
            "Uptime dapat bobot tertinggi (0.25) karena server yang sering down "
            "bakal bikin pengalaman main jadi ga enak. Harga dan RAM sama-sama "
            "penting (0.20) — biaya harus terjangkau tapi spek harus oke.",
            18, GRAY)

# ============================================================
# SLIDE 6 — SEKILAS TOPSIS
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide6, Inches(0), Inches(0), Inches(0.15), H, BLUE)

add_textbox(slide6, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
            "SEKILAS METODE TOPSIS", 36, BLUE, bold=True)

add_textbox(slide6, Inches(1), Inches(1.3), Inches(11), Inches(1),
            "TOPSIS = Technique for Order Preference by Similarity to Ideal Solution", 20, DARK, bold=True)

add_textbox(slide6, Inches(1), Inches(2.3), Inches(11), Inches(4.5),
            "Gampangnya: TOPSIS milih alternatif terbaik berdasarkan "
            "seberapa dekat dia sama solusi ideal.\n\n"
            "Biar lebih gampang — anggap aja kayak lomba lari:\n\n"
            "- Solusi Ideal Positif (A+) = pelari tercepat yang jadi acuan\n"
            "- Solusi Ideal Negatif (A-) = pelari paling lambat (males banget)\n"
            "- Alternatif bagus = yang jaraknya paling dekat ke pelari tercepat,\n"
            "  dan paling jauh dari pelari paling lambat\n\n"
            "TOPSIS ngitung jarak tiap alternatif ke dua titik itu, "
            "terus ngasih nilai preferensi. Makin tinggi nilainya (mendekati 1), "
            "makin bagus alternatifnya.",
            22, DARK)

# ============================================================
# SLIDE 7 — LANGKAH PERHITUNGAN
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide7, Inches(0), Inches(0), Inches(0.15), H, BLUE)

add_textbox(slide7, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
            "LANGKAH PERHITUNGAN", 36, BLUE, bold=True)

steps = [
    ("1. Normalisasi", "Menyamakan skala semua kriteria biar adil"),
    ("2. Normalisasi Terbobot", "Kaliin hasil normalisasi dengan bobot masing-masing kriteria"),
    ("3. Solusi Ideal", "Cari nilai terbaik (A+) dan terburuk (A-) tiap kriteria"),
    ("4. Jarak Euclidean", "Hitung jarak tiap alternatif ke A+ dan A-"),
    ("5. Preferensi", "Hitung nilai V = D- / (D+ + D-), makin tinggi makin baik"),
    ("6. Ranking", "Urutkan dari V tertinggi ke terendah"),
]

tbl = make_table(slide7, len(steps) + 1, 3, Inches(1), Inches(1.3), Inches(11), Inches(4))
style_header(tbl, ["Langkah", "Nama", "Penjelasan"], font_size=14)

for i, (nm, desc) in enumerate(steps):
    style_cell(tbl, i + 1, 0, f"Step {i+1}", font_size=13, align=PP_ALIGN.CENTER, fill=LIGHT_BLUE)
    style_cell(tbl, i + 1, 1, nm, font_size=13, bold=True)
    style_cell(tbl, i + 1, 2, desc, font_size=12)

add_textbox(slide7, Inches(1), Inches(5.8), Inches(11), Inches(1.5),
            "Detail perhitungan lengkap ada di laporan DOCX.",
            16, GRAY)

# ============================================================
# SLIDE 8 — HASIL RANKING TOP 10
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide8, Inches(0), Inches(0), Inches(0.15), H, BLUE)

add_textbox(slide8, Inches(1), Inches(0.3), Inches(10), Inches(0.8),
            "HASIL RANKING (TOP 10)", 36, BLUE, bold=True)

top10 = df_sorted[["Rank", "Platform", "Tier", "Harga (USD/mo)", "RAM (GB)", "Storage (GB)", "vCPU", "V"]].head(10)

tbl = make_table(slide8, 11, 8, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8))
style_header(tbl, ["Rank", "Platform", "Tier", "Harga", "RAM", "Storage", "vCPU", "Nilai V"], font_size=12)

for i in range(10):
    r = top10.iloc[i]
    gold_fill = RGBColor(0xFF, 0xF3, 0xCD) if i == 0 else None
    style_cell(tbl, i + 1, 0, str(int(r["Rank"])), font_size=12, bold=(i == 0),
               align=PP_ALIGN.CENTER, fill=gold_fill)
    style_cell(tbl, i + 1, 1, r["Platform"], font_size=12, bold=(i == 0), fill=gold_fill)
    style_cell(tbl, i + 1, 2, r["Tier"], font_size=12, fill=gold_fill)
    style_cell(tbl, i + 1, 3, f"${r['Harga (USD/mo)']:.2f}", font_size=12, align=PP_ALIGN.CENTER, fill=gold_fill)
    style_cell(tbl, i + 1, 4, f"{r['RAM (GB)']:.0f} GB", font_size=12, align=PP_ALIGN.CENTER, fill=gold_fill)
    style_cell(tbl, i + 1, 5, f"{r['Storage (GB)']:.0f} GB", font_size=12, align=PP_ALIGN.CENTER, fill=gold_fill)
    style_cell(tbl, i + 1, 6, f"{r['vCPU']:.0f}", font_size=12, align=PP_ALIGN.CENTER, fill=gold_fill)
    style_cell(tbl, i + 1, 7, f"{r['V']:.6f}", font_size=12, bold=True, align=PP_ALIGN.CENTER, fill=gold_fill)

# ============================================================
# SLIDE 9 — ANALISIS
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide9, Inches(0), Inches(0), Inches(0.15), H, BLUE)

add_textbox(slide9, Inches(1), Inches(0.3), Inches(8), Inches(0.8),
            "ANALISIS HASIL", 36, BLUE, bold=True)

best = top10.iloc[0]

add_textbox(slide9, Inches(1), Inches(1.3), Inches(11), Inches(2.5),
            f"Juara: {best['Platform']} — {best['Tier']}\n"
            f"Dengan nilai V = {best['V']:.6f} (jauh di atas pesaing)\n\n"
            "Hostinger unggul telak karena:\n"
            "  Harga per GB paling murah — $27.99 dapet 32 GB RAM + 400 GB storage\n"
            "  Dibanding Godlike.host (posisi 2) yang $64.99 untuk spek mirip\n"
            "  Hostinger 2x lebih murah dengan nilai hampir sama!",
            22, DARK)

add_textbox(slide9, Inches(1), Inches(4.3), Inches(11), Inches(1.5),
            "Menariknya, tier 4-16 GB Hostinger (Panel4, Panel2, Panel1) "
            "juga masuk top 10. Ini karena harga promo mereka sangat "
            "agresif dibanding kompetitor dengan spek setara.\n\n"
            "Hosting gratis (Aternos, Minefort, dll) ada di peringkat "
            "bawah karena RAM kecil (1-2 GB) dan uptime ga terjamin.",
            18, GRAY)

# Highlight box
add_shape(slide9, Inches(1), Inches(6), Inches(11), Inches(1), LIGHT_BLUE)
add_textbox(slide9, Inches(1.3), Inches(6.1), Inches(10.5), Inches(0.8),
            "Catatan: Hostinger menawarkan promo diskon besar untuk pembelian tahunan. "
            "Harga reguler mungkin berbeda. Selalu cek website resmi untuk harga terkini.",
            16, GRAY)

# ============================================================
# SLIDE 10 — KESIMPULAN
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, BLUE)

add_textbox(slide10, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
            "KESIMPULAN", 44, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide10, Inches(4.5), Inches(2.5), Inches(4.3), Inches(0.04), WHITE)

add_textbox(slide10, Inches(1), Inches(3), Inches(11), Inches(3),
            f"Berdasarkan perhitungan TOPSIS pada 72 alternatif "
            f"dari 22 platform hosting, alternatif terbaik adalah\n\n"
            f"{best['Platform']} — {best['Tier']}\n"
            f"dengan nilai preferensi V = {best['V']:.6f}\n\n"
            f"Hostinger menawarkan nilai terbaik dengan harga "
            f"${best['Harga (USD/mo)']:.2f}/bulan untuk "
            f"{best['RAM (GB)']:.0f} GB RAM dan "
            f"{best['Storage (GB)']:.0f} GB storage.",
            24, WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide10, Inches(2), Inches(6.3), Inches(9), Inches(1),
            "Terima Kasih  |  Sistem Pendukung Keputusan  |  2026",
            16, RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)

# ============================================================
# SIMPAN
# ============================================================
import os, time
out_path = "Dokumen/SPK_TOPSIS_Minecraft_Hosting.pptx"
# Delete old file with retry
for retry in range(3):
    try:
        if os.path.exists(out_path):
            os.remove(out_path)
        break
    except PermissionError:
        time.sleep(0.5)
prs.save(out_path)
print(f"[OK] PPT: {out_path}")
print(f"     Total slide: {len(prs.slides)}")
